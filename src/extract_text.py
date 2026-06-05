import os

from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document


def load_pdf_text(pdf_path):
    return PyPDFLoader(pdf_path).load()


def load_docx_text(docx_path):

    from docx import Document as DocxDocument

    doc = DocxDocument(docx_path)

    # docx has no real pages, so the whole document is one text block
    text = "\n".join(
        p.text for p in doc.paragraphs if p.text.strip()
    )

    return [
        Document(
            page_content=text,
            metadata={"source": docx_path, "page": 0}
        )
    ]


def load_pptx_text(pptx_path):

    from pptx import Presentation

    prs = Presentation(pptx_path)

    docs = []

    # one Document per slide (slide number acts as the page)
    for slide_number, slide in enumerate(prs.slides):

        lines = []

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                lines.append(shape.text)

        if not lines:
            continue

        docs.append(
            Document(
                page_content="\n".join(lines),
                metadata={"source": pptx_path, "page": slide_number}
            )
        )

    return docs


def load_text_docs(file_path):

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return load_pdf_text(file_path)

    if ext == ".docx":
        return load_docx_text(file_path)

    if ext == ".pptx":
        return load_pptx_text(file_path)

    raise ValueError(f"Unsupported file type: {ext}")
