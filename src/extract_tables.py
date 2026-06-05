import pdfplumber
import pandas as pd
import re
import sqlite3
import os


def clean_column_name(col):

    col = "" if col is None else str(col)

    col = col.replace("\n", " ")

    col = col.strip()

    # remove non-ascii characters
    col = col.encode("ascii", "ignore").decode()

    # remove special characters
    col = re.sub(r"[^\w\s]", "", col)

    # replace spaces with underscores
    col = col.replace(" ", "_")
    col = col.strip("_")

    return col


def clean_cell(x):

    if isinstance(x, str):

        x = x.replace("\n", " ")

        x = " ".join(x.split())

        # remove non-ascii characters
        x = x.encode("ascii", "ignore").decode()

    return x


def build_table_df(rows):
    """Take raw table rows (first row = header) and return a cleaned DataFrame."""

    df = pd.DataFrame(
        rows[1:],
        columns=rows[0]
    )

    # clean column names
    df.columns = [
        clean_column_name(col)
        for col in df.columns
    ]

    # ensure unique, non-empty column names
    seen = {}
    unique_cols = []
    for i, col in enumerate(df.columns):
        col = col or f"col_{i}"
        if col in seen:
            seen[col] += 1
            col = f"{col}_{seen[col]}"
        else:
            seen[col] = 0
        unique_cols.append(col)
    df.columns = unique_cols

    # clean cell values
    df = df.map(clean_cell)

    # auto-detect numeric columns: strip commas/currency,
    # try to_numeric, and keep the conversion if most values parse.
    for col in df.columns:

        stripped = (
            df[col]
            .astype(str)
            .str.replace(",", "", regex=False)
            .str.replace(r"[^\d\.\-eE]", "", regex=True)
        )

        converted = pd.to_numeric(stripped, errors="coerce")

        non_empty = df[col].astype(str).str.strip().ne("").sum()

        if non_empty == 0:
            continue

        parse_rate = converted.notna().sum() / non_empty

        if parse_rate >= 0.8:
            df[col] = converted

    return df


def tables_to_db(tables, db_path):
    """Write a list of {page, table_index, dataframe} into a fresh sqlite db."""

    db_dir = os.path.dirname(db_path)
    if db_dir:
        os.makedirs(db_dir, exist_ok=True)

    if os.path.exists(db_path):
        os.remove(db_path)

    conn = sqlite3.connect(db_path)

    table_metadata = []

    for table_info in tables:

        page = table_info["page"]

        table_index = table_info["table_index"]

        df = table_info["dataframe"]

        table_name = f"page_{page}_table_{table_index}"

        df.to_sql(
            table_name,
            conn,
            if_exists="replace",
            index=False
        )

        table_metadata.append({
            "table_name": table_name,
            "page": page,
            "columns": list(df.columns)
        })

    conn.close()

    return {
        "table_metadata": table_metadata,
        "db_path": db_path
    }


def _resolve_db_path(file_path, db_path):
    if db_path is None:
        name = os.path.splitext(
            os.path.basename(file_path)
        )[0]
        db_path = f"{name}.db"
    return db_path


def extract_pdf_tables(pdf_path, db_path=None):

    db_path = _resolve_db_path(pdf_path, db_path)

    tables = []

    with pdfplumber.open(pdf_path) as pdf:

        for page_number, page in enumerate(pdf.pages):

            extracted_tables = page.extract_tables()

            for table_index, table in enumerate(extracted_tables):

                if not table:
                    continue

                tables.append({
                    "page": page_number,
                    "table_index": table_index,
                    "dataframe": build_table_df(table)
                })

    return tables_to_db(tables, db_path)


def extract_docx_tables(docx_path, db_path=None):

    from docx import Document as DocxDocument

    db_path = _resolve_db_path(docx_path, db_path)

    doc = DocxDocument(docx_path)

    tables = []

    for table_index, table in enumerate(doc.tables):

        rows = [
            [cell.text for cell in row.cells]
            for row in table.rows
        ]

        if not rows:
            continue

        tables.append({
            "page": 0,
            "table_index": table_index,
            "dataframe": build_table_df(rows)
        })

    return tables_to_db(tables, db_path)


def extract_pptx_tables(pptx_path, db_path=None):

    from pptx import Presentation

    db_path = _resolve_db_path(pptx_path, db_path)

    prs = Presentation(pptx_path)

    tables = []

    for slide_number, slide in enumerate(prs.slides):

        table_index = 0

        for shape in slide.shapes:

            if not shape.has_table:
                continue

            rows = [
                [cell.text for cell in row.cells]
                for row in shape.table.rows
            ]

            if not rows:
                continue

            tables.append({
                "page": slide_number,
                "table_index": table_index,
                "dataframe": build_table_df(rows)
            })

            table_index += 1

    return tables_to_db(tables, db_path)


def extract_tables(file_path, db_path=None):

    ext = os.path.splitext(file_path)[1].lower()

    if ext == ".pdf":
        return extract_pdf_tables(file_path, db_path)

    if ext == ".docx":
        return extract_docx_tables(file_path, db_path)

    if ext == ".pptx":
        return extract_pptx_tables(file_path, db_path)

    raise ValueError(f"Unsupported file type: {ext}")
